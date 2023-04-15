import cv2  # Importation de la bibliothèque OpenCV pour la capture vidéo et le traitement des images
import mediapipe as mp  # Importation de la bibliothèque Mediapipe pour la détection des mains
import numpy as np  # Importation de la bibliothèque NumPy pour les opérations mathématiques
from mediapipe.framework.formats import landmark_pb2  # Importation de la bibliothèque landmark_pb2 pour le traitement des points de repère
import time  # Importation de la bibliothèque time pour la mesure du temps
import win32api  # Importation de la bibliothèque win32api pour l'interaction avec le système d'exploitation Windows
import pyautogui  # Importation de la bibliothèque pyautogui pour la simulation des actions de souris
from pptx import Presentation  # Importation de la bibliothèque pptx pour la manipulation des présentations PowerPoint

mp_drawing = mp.solutions.drawing_utils  # Initialisation de la bibliothèque Mediapipe pour dessiner les points de repère sur les mains
mp_hands = mp.solutions.hands  # Initialisation de la bibliothèque Mediapipe pour la détection des mains
click = 0  # Initialisation de la variable click à zéro 
open_hand = False  # Initialisation de la variable open_hand à False
closed_hand = False  # Initialisation de la variable closed_hand à False

video = cv2.VideoCapture(0)  # Initialisation de la capture vidéo à partir de la webcam

with mp_hands.Hands(min_detection_confidence=0.8, min_tracking_confidence=0.8) as hands:  # Définition des paramètres de détection des mains
    prs = Presentation('my_presentation.pptx')  # Chargement du fichier PowerPoint à contrôler
    slide_count = len(prs.slides)  # Calcul du nombre de diapositives dans la présentation
    current_slide = 0  # Initialisation de la variable current_slide à zéro

    while video.isOpened():  # Boucle principale de traitement de la vidéo
        _, frame = video.read()  # Lecture d'une image à partir de la webcam
        image = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)  # Conversion de l'image en couleur

        image = cv2.flip(image, 1)  # Renversement horizontal de l'image

        imageHeight, imageWidth, _ = image.shape  # Récupération des dimensions de l'image

        results = hands.process(image)  # Détection des mains dans l'image à l'aide de Mediapipe

        image = cv2.cvtColor(image, cv2.COLOR_RGB2BGR)  # Conversion de l'image en BGR

        if results.multi_hand_landmarks:  # Si des mains sont détectées dans l'image
            for num, hand in enumerate(results.multi_hand_landmarks):  # Pour chaque main détectée
                mp_drawing.draw_landmarks(image, hand, mp_hands.HAND_CONNECTIONS,  # Dessin des points de repère sur la main
                                          mp_drawing.DrawingSpec(color=(250, 44, 250), thickness=2, circle_radius=2),
                                          )

            # Récupération des positions des bouts de l'index et du pouce pour calculer la distance entre eux
            for handLandmarks in results.multiresults.multi_hand_landmarks:
                        # Mise à jour de la diapositive en fonction du geste de la main
                if closed_hand: # Si la main est fermée
                    if current_slide < slide_count - 1: # Si la diapositive actuelle n'est pas la dernière
                        current_slide += 1 # Passer à la diapositive suivante
                        prs.slides[current_slide].select() # Sélectionner la diapositive dans PowerPoint
                elif open_hand: # Si la main est ouverte
                    if current_slide > 0: # Si la diapositive actuelle n'est pas la première
                        current_slide -= 1 # Revenir à la diapositive précédente
                        prs.slides[current_slide].select() # Sélectionner la diapositive dans PowerPoint

                # Affichage de l'image avec les landmarks dessinés
                cv2.imshow('Hand Tracking', image)

                if cv2.waitKey(10) & 0xFF == ord('q'): # Si la touche "q" est pressée, arrêter le programme
                    break

# Fermer la fenêtre de la webcam et sauvegarder la présentation
video.release()
cv2.destroyAllWindows()
prs.save('my_presentation.pptx')

